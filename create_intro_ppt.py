from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

prs = Presentation()

# Theme-ish colors
BG = RGBColor(247, 243, 232)
TITLE = RGBColor(20, 40, 80)
ACCENT = RGBColor(214, 93, 14)
BODY = RGBColor(35, 35, 35)


def style_title(shape):
    tf = shape.text_frame
    for p in tf.paragraphs:
        for r in p.runs:
            r.font.name = "Calibri"
            r.font.bold = True
            r.font.size = Pt(40)
            r.font.color.rgb = TITLE


def add_title_slide(title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG

    slide.shapes.title.text = title
    style_title(slide.shapes.title)

    sub = slide.placeholders[1]
    sub.text = subtitle
    p = sub.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    for r in p.runs:
        r.font.name = "Calibri"
        r.font.size = Pt(22)
        r.font.color.rgb = ACCENT


def add_bullets_slide(title, bullets, note=None):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG

    slide.shapes.title.text = title
    style_title(slide.shapes.title)

    tf = slide.placeholders[1].text_frame
    tf.clear()
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = b
        p.level = 0
        p.space_after = Pt(8)
        for r in p.runs:
            r.font.name = "Calibri"
            r.font.size = Pt(24)
            r.font.color.rgb = BODY

    if note:
        box = slide.shapes.add_textbox(Inches(0.6), Inches(6.4), Inches(12.2), Inches(0.6))
        t = box.text_frame
        t.text = note
        p = t.paragraphs[0]
        p.alignment = PP_ALIGN.RIGHT
        for r in p.runs:
            r.font.name = "Calibri"
            r.font.italic = True
            r.font.size = Pt(14)
            r.font.color.rgb = ACCENT


add_title_slide(
    "MNIST ANN: How a Tiny Brain Learns Digits",
    "Fast, funny, and fully NumPy (no deep-learning frameworks harmed in this experiment)",
)

add_bullets_slide(
    "Why This Matters",
    [
        "A neural net is just matrix math + stubborn optimism.",
        "We control every lever: inputs, hidden nodes, weights, learning rate, epochs.",
        "Because we built it from scratch, we can explain every prediction.",
    ],
    note="Yes, this is legal machine learning with only NumPy.",
)

add_bullets_slide(
    "The Model in One Picture",
    [
        "Input: 784 pixels (28x28) per image.",
        "Hidden layer: learns useful patterns (strokes, corners, curves).",
        "Output: 10 probabilities for digits 0-9.",
    ],
    note="Think: compressed detective work from raw pixels to class probabilities.",
)

add_bullets_slide(
    "Forward Pass (Prediction Time)",
    [
        "Step 1: Weighted sum + bias in hidden layer.",
        "Step 2: ReLU keeps positive signals, zeros out weak negatives.",
        "Step 3: Softmax turns scores into probabilities.",
        "Result: The model picks the highest-probability digit.",
    ],
    note="Forward pass = the network making its best guess.",
)

add_bullets_slide(
    "Backpropagation (Learning Time)",
    [
        "Compare prediction with true label to get error.",
        "Send error backward through the layers (chain rule).",
        "Compute gradients for each weight and bias.",
        "Update parameters: old value - learning_rate x gradient.",
    ],
    note="Backprop = polite but relentless correction.",
)

add_bullets_slide(
    "The Speed/Accuracy Trade-Off",
    [
        "More hidden nodes: usually better accuracy, slower epochs.",
        "Bigger learning rate: faster progress, but can overshoot.",
        "More epochs: better fit, until overfitting knocks.",
        "Batch size changes throughput and generalization behavior.",
    ],
    note="Training is optimization, not wishful thinking.",
)

add_bullets_slide(
    "What We Experimented With",
    [
        "Hidden nodes sweep: 32, 64, 128, 256.",
        "Initialization scale sweep: 0.25, 0.5, 1.0, 2.0.",
        "Speed sensitivity: batch size x learning rate grid.",
        "Learning curves tracked every epoch.",
    ],
    note="We don't just report one score; we map the behavior landscape.",
)

add_bullets_slide(
    "How to Read the Curves",
    [
        "Train and validation accuracy rising together: healthy learning.",
        "Train up, validation flat/down: overfitting warning.",
        "Both low and flat: underfitting or poor hyperparameters.",
        "Epoch-time curve reveals computational cost directly.",
    ],
    note="Curves are the model's diary: confidence, confusion, and growth.",
)

add_bullets_slide(
    "Before the Live Demo",
    [
        "We know what each knob does.",
        "We know where speed is gained and where accuracy is lost.",
        "We can justify architecture choices with evidence.",
        "Now let's demonstrate training and predictions live.",
    ],
    note="Time to make the digits confess.",
)

out_path = "/home/bantu/Documents/MNIST_ANN/MNIST_ANN_Intro_Presentation.pptx"
prs.save(out_path)
print(out_path)
